import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# フォント設定
plt.rcParams['font.size'] = 18
plt.rcParams['font.family'] = 'Helvetica'
plt.rcParams['font.weight'] = 'bold'

# CSVファイルの読み込み
csv_file = '230722-240126_murata.csv'
df = pd.read_csv(csv_file, parse_dates=['time'])

# CSVファイル名から拡張子を除いた基本名を取得
base_name = os.path.splitext(csv_file)[0]

# フィルタリング
start_date = '2023-7-22'
end_date = '2024-1-31'
filtered_df = df[(df['time'] >= start_date) & (df['time'] <= end_date)]

# グラフの作成
fig2, ax2 = plt.subplots(figsize=(22, 9))

# 0-6時と6-24時のデータをフィルタリング
morning_df = filtered_df[filtered_df['time'].dt.hour < 6]
daytime_df = filtered_df[filtered_df['time'].dt.hour >= 6]

# コード2に基づくグラフの作成
fig2, ax2 = plt.subplots(figsize=(22, 9))
morning_df = filtered_df[filtered_df['time'].dt.hour < 6]
morning_stats = morning_df.groupby(morning_df['time'].dt.date)['glucose'].agg(['mean', 'std']).reset_index()
daytime_df = filtered_df[filtered_df['time'].dt.hour >= 6]
daytime_stats = daytime_df.groupby(daytime_df['time'].dt.date)['glucose'].agg(['mean', 'std']).reset_index()
morning_stats.rename(columns={'time': 'Date'}, inplace=True)
daytime_stats.rename(columns={'time': 'Date'}, inplace=True)
morning_avg = morning_df['glucose'].mean()
daytime_avg = daytime_df['glucose'].mean()
morning_dates = mdates.date2num(morning_stats['Date'])
daytime_dates = mdates.date2num(daytime_stats['Date'])

# エラーバーと平均値線のプロット設定の変更
marker_size = 12  # マーカーのサイズ
line_width = 2    # 線の太さ

ax2.errorbar(morning_dates, morning_stats['mean'], yerr=morning_stats['std'], fmt='o-', color='blue', label='0-6 h', capsize=5, markersize=marker_size, linewidth=line_width)
ax2.errorbar(daytime_dates, daytime_stats['mean'], yerr=daytime_stats['std'], fmt='o-', color='red', label='6-24 h', capsize=5, markersize=marker_size, linewidth=line_width)
ax2.axhline(y=morning_avg, color='blue', linestyle='--', label='Average 0-6 h', linewidth=line_width)
ax2.axhline(y=daytime_avg, color='lightcoral', linestyle='--', label='Average 6-24 h', linewidth=line_width)

ax2.xaxis.set_major_locator(mdates.DayLocator())
ax2.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
ax2.set_ylim(40, 140)
ax2.set_xlabel('Date', fontweight='bold', fontsize=12)
ax2.set_ylabel('Average glucose level / mg/dL', fontweight='bold', fontsize=24)
plt.xticks(rotation=45)

# x軸の日付が切れないように調整
plt.gcf().subplots_adjust(bottom=0.2)


# グラフ2の保存 (ファイル名を '211107_satoh-2.png' に変更、背景を透明に)
fig2.savefig(f'{base_name}-2.png', transparent=True)

# グラフを表示
plt.show()

# 平均値の計算
morning_avg = morning_df['glucose'].mean()
daytime_avg = daytime_df['glucose'].mean()


# 凡例を表示させない
ax2.legend().set_visible(False)

# グラフ2の保存（背景を透明にし、凡例を表示せず）
fig2.savefig(f'{base_name}-2.png', transparent=True)


# 凡例を別に保存
fig_legend, ax_legend = plt.subplots(figsize=(3, 2))  # 凡例の図のサイズを設定（縦長に）
handles, labels = ax2.get_legend_handles_labels()
leg = ax_legend.legend(handles, labels, loc='center', ncol=1, frameon=False)  # ncol=1 で縦一列に設定
for text in leg.get_texts():
    text.set_fontsize(12)  # 凡例のテキストのフォントサイズを設定
ax_legend.axis('off')
legend_path = f'{base_name}-legend.png'
fig_legend.tight_layout()
fig_legend.savefig(legend_path, transparent=True)


# PowerPointに挿入するためのテキストボックスを準備
textbox_text = f"Average 0-6 h: {morning_avg:.1f} mg/dL\nAverage 6-24 h: {daytime_avg:.1f} mg/dL"

# PowerPointファイルに挿入
prs = Presentation('presentation_a4_background_yoko.pptx')

# 既存の最初のスライドを取得
slide = prs.slides[0]

# グラフ、凡例、平均値テキスト、測定期間テキストの挿入位置を設定
graph_left = Inches(-2 / 2.54)  # -2 cm
graph_top = Inches(3 / 2.54)    # 2 cm
graph_height = Inches(5)

legend_left = Inches(20.5 / 2.54)  # 16 cm
legend_top = Inches(14.5 / 2.54)    # 1 cm

textbox_left = Inches(1 / 2.54)  # 1 cm
textbox_top = Inches(16.5 / 2.54)  # 14 cm
textbox_width = Inches(5)
textbox_height = Inches(1)

textbox2_left = Inches(21.2 / 2.54)   # 1 cm
textbox2_top = Inches(-0.05 / 2.54)   # 16 cm
textbox2_width = Inches(6)
textbox2_height = Inches(1)

# グラフ、凡例、平均値テキストを挿入
slide.shapes.add_picture(f'{base_name}-2.png', graph_left, graph_top, height=graph_height)
slide.shapes.add_picture(legend_path, legend_left, legend_top)

# 平均値テキストボックスを挿入
textbox = slide.shapes.add_textbox(textbox_left, textbox_top, width=textbox_width, height=textbox_height)
p = textbox.text_frame.add_paragraph()
p.text = f"Average 0-6 h: {morning_avg:.1f} mg/dL\nAverage 6-24 h: {daytime_avg:.1f} mg/dL"
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Helvetica'
p.alignment = PP_ALIGN.LEFT

# 測定期間テキストボックスを挿入
textbox2 = slide.shapes.add_textbox(textbox2_left, textbox2_top, width=textbox2_width, height=textbox2_height)
p2 = textbox2.text_frame.add_paragraph()
p2.text = f"{start_date} 〜 {end_date}"
p2.font.size = Pt(14)
p2.font.bold = True
p2.font.name = 'Helvetica'
p2.alignment = PP_ALIGN.LEFT

# PowerPointファイルの保存
prs.save('updated_presentation_yoko.pptx')