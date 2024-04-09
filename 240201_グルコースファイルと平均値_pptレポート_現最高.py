import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
import matplotlib.dates as md
from datetime import datetime, timedelta
from scipy.interpolate import make_interp_spline
import matplotlib.patheffects as PathEffects
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# フォント設定
plt.rcParams['font.size'] = 24
plt.rcParams['font.family'] = 'Helvetica'
plt.rcParams['font.weight'] = 'bold'

# CSVファイルの読み込み
csv_file = '231112-1226_tateno.csv'
df = pd.read_csv(csv_file, parse_dates=['time'])


# CSVファイル名から拡張子を除いた基本名を取得
base_name = os.path.splitext(csv_file)[0]


colors = [
    '#000000', # 黒
    '#FFFF00', # 黄
    '#0000FF', # 青
    '#FF0000', # 赤
    '#87CEFA', # ライトスカイブルー
    '#FFD700', # ゴールド
    '#808080', # 灰色
    '#20B2AA', # ライトシーグリーン
    '#FF1493', # 濃いピンク
    '#000080', # ネイビー
    '#FFA500', # オレンジ
    '#00FFFF', # シアン
    '#C71585', # ミディアムバイオレットレッド
    '#800080', # 紫
    '#D2691E', # チョコレート色
    '#FF00FF', # マゼンタ
    '#C0C0C0', # シルバー
    '#2E8B57', # シーグリーン
    '#FF6347', # トマト色
    '#BA55D3', # ミディアムオーキッド
    '#FF4500', # オレンジレッド
    '#8A2BE2', # ブルーバイオレット
    '#008000', # 緑
    '#8B4513', # サドルブラウン
    '#32CD32', # ライムグリーン
    '#40E0D0', # ターコイズ
    '#5F9EA0', # ケイデットブルー
    '#008080', # ティール
    '#DA70D6', # オーキッド
    '#228B22', # フォレストグリーン
    '#A0522D', # シエナ
    '#CD5C5C', # インディアンレッド
    '#FFC0CB', # ピンク
    '#FA8072', # サーモン
    '#B0C4DE', # ライトスチールブルー
    '#ADD8E6', # ライトブルー
    '#DEB887', # バーリウッド
    '#F5DEB3', # ホウィート
    '#FFFACD', # レモンシフォン
    '#E0FFFF' # ライトシアン
]



# フィルタリング
start_date = pd.to_datetime('2023-11-12')
end_date = pd.to_datetime('2023-12-01')
filtered_df = df[(df['time'] >= start_date) & (df['time'] <= end_date)]

# Generate a list of all dates within the range and pre-assign them an index
all_dates = pd.date_range(start=start_date, end=end_date, freq='D')
date_to_index = {date.date(): i for i, date in enumerate(all_dates)}

# 確認：colors リストに十分な色が含まれていることを確認
assert len(colors) >= len(all_dates), "There are not enough colors for each date."

# 確認：date_to_index が全ての日付を含んでいることを確認
assert all(date in date_to_index for date in filtered_df['time'].dt.date.unique()), "Not all dates are in date_to_index mapping."


# グラフの作成
fig1, ax1 = plt.subplots(figsize=(28, 12), constrained_layout=True)
ax1.set_xlabel("Time", fontweight='bold', fontsize=24)
ax1.set_ylabel("Interstitial glucose level / mg/dL", fontweight='bold', fontsize=24)

# グラフの設定
base_date = datetime(1900, 1, 1)
next_day = base_date + timedelta(days=1)
ax1.set_xlim([mdates.date2num(base_date), mdates.date2num(next_day)])




# しきい値の設定
time_threshold = timedelta(minutes=30)
base_date = datetime(1900, 1, 1)



# 凡例の色とラベルを保持するための辞書を作成
legend_labels = {}

# 日付ごとにデータを処理
for date, gdf in filtered_df.groupby(filtered_df['time'].dt.date):
    # 正しい color_index を取得
    color_index = date_to_index[date] % len(colors)
    color = colors[color_index]
    
    # 凡例用のラベルと色を辞書に保存
    legend_labels[date.strftime('%Y-%m-%d')] = color

    
    gdf = gdf.sort_values('time')
    gdf['time'] = gdf['time'].apply(lambda x: datetime.combine(base_date, x.time()))

    # セグメントを初期化
    segments = []
    current_segment = [gdf.iloc[0]]

    # セグメントの作成
    for j in range(1, len(gdf)):
        if (gdf.iloc[j]['time'] - gdf.iloc[j - 1]['time']) > time_threshold:
            segments.append(current_segment)
            current_segment = [gdf.iloc[j]]
        else:
            current_segment.append(gdf.iloc[j])
    segments.append(current_segment)

    # 各セグメントに対してスプライン曲線を描画
    for segment in segments:
        if len(segment) < 2:
            continue
        segment_df = pd.DataFrame(segment)
        x = mdates.date2num(segment_df['time'])
        y = segment_df['glucose']

        # 重複するx値の処理
        if len(np.unique(x)) < len(x):
            unique_x = np.unique(x)
            y_means = [y[x == ux].mean() for ux in unique_x]
            x, y = unique_x, y_means


        # スプライン曲線の作成と描画
        spline = make_interp_spline(x, y, k=1)
        x_new = np.linspace(x.min(), x.max(), 300)
        y_smooth = spline(x_new)
        line, = ax1.plot(mdates.num2date(x_new), y_smooth, color=color, linewidth=4)
        # 影のようなエフェクトを線に適用する
        effect = PathEffects.withStroke(linewidth=5, foreground='black')
        line.set_path_effects([effect])

# グラフの設定
ax1.set_xlim([mdates.date2num(base_date), mdates.date2num(base_date + timedelta(days=1))])
ax1.set_ylim(50, 200)
ax1.xaxis.set_major_locator(mdates.HourLocator(interval=1))
ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M'))
ax1.grid(which='major', linestyle='--', linewidth=0.5)
ax1.legend([date.strftime('%Y-%m-%d') for date in filtered_df['time'].dt.date.unique()], loc='upper left', bbox_to_anchor=(1, 1))
plt.xticks(rotation=45)




# 凡例をグラフの外に配置し、フォントサイズを調整する
# 凡例のラベルと色を使用して凡例を作成
handles = [plt.Line2D([0], [0], color=color, linewidth=4) for color in legend_labels.values()]
leg = ax1.legend(
    handles,
    legend_labels.keys(),
    loc='upper left', 
    bbox_to_anchor=(1.06, 1),
    prop={'size': 18},  # 凡例のフォントサイズを小さくする
    title='Date', 
    title_fontsize='18'  # 凡例のタイトルのフォントサイズを小さくする
)

# 凡例の行間を詰める
leg.get_frame().set_linewidth(0.0)  # 凡例の枠線を消す

# tight_layoutを呼び出して、グラフのレイアウトを自動的に調整する
fig1.tight_layout()

# 凡例を含むグラフを保存する際に余白を調整する
# 注意: bbox_inches='tight' を使用していると、tight_layoutの調整が無視される可能性があるため、適切なパディングを指定する
fig1.savefig(f'{base_name}-1.png', bbox_extra_artists=(leg,), bbox_inches='tight', pad_inches=0.5, transparent=True)




# グラフの作成
fig2, ax2 = plt.subplots(figsize=(22, 9), constrained_layout=True)

# 0-6時と6-24時のデータをフィルタリング
morning_df = filtered_df[filtered_df['time'].dt.hour < 6]
daytime_df = filtered_df[filtered_df['time'].dt.hour >= 6]

# コード2に基づくグラフの作成
fig2, ax2 = plt.subplots(figsize=(22, 9))
morning_df = filtered_df[filtered_df['time'].dt.hour < 6]
morning_stats = morning_df.groupby(morning_df['time'].dt.date)['glucose'].agg(['mean', 'std']).reset_index()
daytime_df = filtered_df[filtered_df['time'].dt.hour >= 6]
daytime_stats = daytime_df.groupby(daytime_df['time'].dt.date)['glucose'].agg(['mean', 'std']).reset_index()
morning_stats.rename(columns={'time': 'Date'}, inplace=True)
daytime_stats.rename(columns={'time': 'Date'}, inplace=True)
morning_avg = morning_df['glucose'].mean()
daytime_avg = daytime_df['glucose'].mean()
morning_dates = mdates.date2num(morning_stats['Date'])
daytime_dates = mdates.date2num(daytime_stats['Date'])

# エラーバーと平均値線のプロット設定の変更
marker_size = 12  # マーカーのサイズ
line_width = 2    # 線の太さ

ax2.errorbar(morning_dates, morning_stats['mean'], yerr=morning_stats['std'], fmt='o-', color='blue', label='0-6 h', capsize=5, markersize=marker_size, linewidth=line_width)
ax2.errorbar(daytime_dates, daytime_stats['mean'], yerr=daytime_stats['std'], fmt='o-', color='red', label='6-24 h', capsize=5, markersize=marker_size, linewidth=line_width)
ax2.axhline(y=morning_avg, color='blue', linestyle='--', label='Average 0-6 h', linewidth=line_width)
ax2.axhline(y=daytime_avg, color='lightcoral', linestyle='--', label='Average 6-24 h', linewidth=line_width)

ax2.xaxis.set_major_locator(mdates.DayLocator())
ax2.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))
ax2.set_ylim(50, 160)
ax2.set_xlabel('Date', fontweight='bold', fontsize=12)
ax2.set_ylabel('Average glucose level / mg/dL', fontweight='bold', fontsize=24)
ax2.legend()
plt.xticks(rotation=45)

# x軸の日付が切れないように調整
plt.gcf().subplots_adjust(bottom=0.2)

# 凡例を含めてグラフ2を保存
ax2.legend(loc='upper right', bbox_to_anchor=(1.21, 0.95), borderaxespad=0.)
fig2.savefig(f'{base_name}-2.png', bbox_inches='tight', transparent=True)


# グラフを表示
plt.show()

# 平均値の計算
morning_avg = morning_df['glucose'].mean()
daytime_avg = daytime_df['glucose'].mean()







# PowerPointファイルに挿入
prs = Presentation('presentation_a4_background-2.pptx')

# 最初のスライドを取得（または新しいスライドを作成）
if len(prs.slides) > 0:
    slide = prs.slides[0]
else:
    slide_layout = prs.slide_layouts[5]  # 適切なレイアウトを選択
    slide = prs.slides.add_slide(slide_layout)

# グラフ1を挿入
slide.shapes.add_picture(f'{base_name}-1.png', Inches(0.8 / 2.54), Inches(5 / 2.54), width=Inches(18 / 2.54))
# グラフ2を挿入
slide.shapes.add_picture(f'{base_name}-2.png', Inches(1.5 / 2.54), Inches(18 / 2.54), width=Inches(16 / 2.54))








# 凡例、平均値テキスト、測定期間テキストの挿入位置を設定
textbox_left = Inches(1 / 2.54)  # 1 cm
textbox_top = Inches(25.4 / 2.54)  # 14 cm
textbox_width = Inches(6 / 2.54)
textbox_height = Inches(1.5 / 2.54)

textbox2_left = Inches(13.5 / 2.54)   #  cm
textbox2_top = Inches(0.04 / 2.54)   #  cm
textbox2_width = Inches(6 / 2.54)
textbox2_height = Inches(1 / 2.54)

textbox3_left = Inches(3.0 / 2.54)   #  cm
textbox3_top = Inches(1.4 / 2.54)   #  cm
textbox3_width = Inches(6 / 2.54)
textbox3_height = Inches(1 / 2.54)

# 平均値テキストボックスを挿入
textbox = slide.shapes.add_textbox(textbox_left, textbox_top, width=textbox_width, height=textbox_height)
p = textbox.text_frame.add_paragraph()
p.text = f"0-6時の平均値: {morning_avg:.1f} mg/dL\n6-24時の平均値: {daytime_avg:.1f} mg/dL"
p.font.size = Pt(12)
p.font.bold = True
p.font.name = 'Helvetica'
p.alignment = PP_ALIGN.LEFT

# 測定期間テキストボックスを挿入
textbox2 = slide.shapes.add_textbox(textbox2_left, textbox2_top, width=textbox2_width, height=textbox2_height)
p2 = textbox2.text_frame.add_paragraph()

# start_dateとend_dateがdatetimeオブジェクトの場合、strftimeメソッドを使用して書式設定
formatted_start_date = start_date.strftime('%Y-%m-%d')
formatted_end_date = end_date.strftime('%Y-%m-%d')
p2.text = f"{formatted_start_date} 〜 {formatted_end_date}"
p2.font.size = Pt(12)
p2.font.bold = True
p2.font.name = 'Helvetica'
p2.alignment = PP_ALIGN.LEFT

# 測定期間テキストボックスを挿入
textbox3 = slide.shapes.add_textbox(textbox3_left, textbox3_top, width=textbox3_width, height=textbox3_height)
p3 = textbox3.text_frame.add_paragraph()
p3.text = f"{base_name}"
p3.font.size = Pt(12)
p3.font.bold = True
p3.font.name = 'Helvetica'
p3.alignment = PP_ALIGN.LEFT



# PowerPointファイルの保存
prs.save('updated_presentation-2.pptx')
